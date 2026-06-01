"""Document conversion module."""

import io
import re
import time
from pathlib import Path
from typing import List, Optional, Union

import chardet
import magic
from bs4 import BeautifulSoup
from markdownify import markdownify as md
from PIL import Image
from rich.console import Console

from documind_ai.core.models import (
    ConversionResult,
    DocumentMetadata,
    DocumentType,
    ProcessingConfig,
)
from documind_ai.utils.file_utils import get_file_extension, sanitize_filename

console = Console()


class DocumentConverter:
    """Convert various document formats to Markdown."""

    # File extension to DocumentType mapping
    EXTENSION_MAP = {
        '.pdf': DocumentType.PDF,
        '.docx': DocumentType.DOCX,
        '.doc': DocumentType.DOC,
        '.xlsx': DocumentType.XLSX,
        '.xls': DocumentType.XLS,
        '.pptx': DocumentType.PPTX,
        '.ppt': DocumentType.PPT,
        '.html': DocumentType.HTML,
        '.htm': DocumentType.HTM,
        '.txt': DocumentType.TXT,
        '.md': DocumentType.MD,
        '.markdown': DocumentType.MARKDOWN,
        '.csv': DocumentType.CSV,
        '.json': DocumentType.JSON,
        '.xml': DocumentType.XML,
        '.epub': DocumentType.EPUB,
        '.mobi': DocumentType.MOBI,
        '.azw': DocumentType.AZW,
        '.azw3': DocumentType.AZW3,
        '.rtf': DocumentType.RTF,
        '.odt': DocumentType.ODT,
        '.ods': DocumentType.ODS,
        '.odp': DocumentType.ODP,
    }

    # Image extensions
    IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.svg'}

    def __init__(self, config: Optional[ProcessingConfig] = None):
        """Initialize converter.
        
        Args:
            config: Processing configuration
        """
        self.config = config or ProcessingConfig()
        self._load_handlers()

    def _load_handlers(self) -> None:
        """Load conversion handlers."""
        self.handlers = {
            DocumentType.PDF: self._convert_pdf,
            DocumentType.DOCX: self._convert_docx,
            DocumentType.DOC: self._convert_doc,
            DocumentType.XLSX: self._convert_xlsx,
            DocumentType.XLS: self._convert_xls,
            DocumentType.PPTX: self._convert_pptx,
            DocumentType.PPT: self._convert_ppt,
            DocumentType.HTML: self._convert_html,
            DocumentType.HTM: self._convert_html,
            DocumentType.TXT: self._convert_text,
            DocumentType.MD: self._convert_markdown,
            DocumentType.MARKDOWN: self._convert_markdown,
            DocumentType.CSV: self._convert_csv,
            DocumentType.JSON: self._convert_json,
            DocumentType.XML: self._convert_xml,
            DocumentType.RTF: self._convert_rtf,
            DocumentType.ODT: self._convert_odt,
            DocumentType.ODS: self._convert_ods,
            DocumentType.ODP: self._convert_odp,
        }

    def detect_document_type(self, file_path: Path) -> DocumentType:
        """Detect document type from file.
        
        Args:
            file_path: Path to file
            
        Returns:
            Detected document type
        """
        # Check by extension first
        ext = get_file_extension(file_path).lower()
        if ext in self.EXTENSION_MAP:
            return self.EXTENSION_MAP[ext]
        
        # Check if image
        if ext in self.IMAGE_EXTENSIONS:
            return DocumentType.IMAGE
        
        # Try MIME type detection
        try:
            mime = magic.from_file(str(file_path), mime=True)
            mime_map = {
                'application/pdf': DocumentType.PDF,
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': DocumentType.DOCX,
                'application/msword': DocumentType.DOC,
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': DocumentType.XLSX,
                'application/vnd.ms-excel': DocumentType.XLS,
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': DocumentType.PPTX,
                'application/vnd.ms-powerpoint': DocumentType.PPT,
                'text/html': DocumentType.HTML,
                'text/plain': DocumentType.TXT,
                'text/markdown': DocumentType.MD,
                'text/csv': DocumentType.CSV,
                'application/json': DocumentType.JSON,
                'application/xml': DocumentType.XML,
                'application/epub+zip': DocumentType.EPUB,
                'application/rtf': DocumentType.RTF,
                'application/vnd.oasis.opendocument.text': DocumentType.ODT,
                'application/vnd.oasis.opendocument.spreadsheet': DocumentType.ODS,
                'application/vnd.oasis.opendocument.presentation': DocumentType.ODP,
            }
            if mime in mime_map:
                return mime_map[mime]
        except Exception:
            pass
        
        return DocumentType.UNKNOWN

    def convert(
        self,
        file_path: Union[str, Path],
        output_path: Optional[Path] = None,
    ) -> ConversionResult:
        """Convert document to Markdown.
        
        Args:
            file_path: Path to input file
            output_path: Optional path for output markdown
            
        Returns:
            Conversion result
        """
        start_time = time.time()
        file_path = Path(file_path)
        errors = []
        warnings = []

        # Validate file exists
        if not file_path.exists():
            return ConversionResult(
                success=False,
                content="",
                document_type=DocumentType.UNKNOWN,
                metadata=DocumentMetadata(),
                file_path=file_path,
                errors=[f"File not found: {file_path}"],
            )

        # Check file size
        file_size = file_path.stat().st_size
        if file_size > self.config.max_file_size:
            return ConversionResult(
                success=False,
                content="",
                document_type=DocumentType.UNKNOWN,
                metadata=DocumentMetadata(file_size=file_size),
                file_path=file_path,
                errors=[f"File too large: {file_size} bytes (max: {self.config.max_file_size})"],
            )

        # Detect document type
        doc_type = self.detect_document_type(file_path)
        
        if doc_type == DocumentType.UNKNOWN:
            return ConversionResult(
                success=False,
                content="",
                document_type=doc_type,
                metadata=DocumentMetadata(file_size=file_size),
                file_path=file_path,
                errors=[f"Unknown document type: {file_path}"],
            )

        # Get appropriate handler
        handler = self.handlers.get(doc_type)
        if not handler:
            return ConversionResult(
                success=False,
                content="",
                document_type=doc_type,
                metadata=DocumentMetadata(file_size=file_size),
                file_path=file_path,
                errors=[f"No handler for document type: {doc_type}"],
            )

        # Perform conversion
        try:
            content, metadata = handler(file_path)
            metadata.file_size = file_size
            
            # Clean up content
            content = self._clean_content(content)
            
            # Save markdown if requested
            if self.config.save_markdown and output_path:
                output_path.parent.mkdir(parents=True, exist_ok=True)
                output_path.write_text(content, encoding='utf-8')
            
            processing_time = time.time() - start_time
            
            return ConversionResult(
                success=True,
                content=content,
                document_type=doc_type,
                metadata=metadata,
                file_path=file_path,
                markdown_path=output_path,
                warnings=warnings,
                processing_time=processing_time,
            )
            
        except Exception as e:
            errors.append(f"Conversion error: {str(e)}")
            return ConversionResult(
                success=False,
                content="",
                document_type=doc_type,
                metadata=DocumentMetadata(file_size=file_size),
                file_path=file_path,
                errors=errors,
                warnings=warnings,
            )

    def _clean_content(self, content: str) -> str:
        """Clean and normalize content.
        
        Args:
            content: Raw content
            
        Returns:
            Cleaned content
        """
        # Remove excessive whitespace
        content = re.sub(r'\n{4,}', '\n\n\n', content)
        
        # Remove trailing whitespace
        content = re.sub(r'[ \t]+$', '', content, flags=re.MULTILINE)
        
        # Normalize line endings
        content = content.replace('\r\n', '\n').replace('\r', '\n')
        
        return content.strip()

    def _convert_pdf(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert PDF to Markdown."""
        import fitz  # PyMuPDF
        
        doc = fitz.open(file_path)
        content_parts = []
        
        metadata = DocumentMetadata(
            title=doc.metadata.get('title'),
            author=doc.metadata.get('author'),
            subject=doc.metadata.get('subject'),
            page_count=len(doc),
        )
        
        # Extract text from each page
        for page_num, page in enumerate(doc, 1):
            if self.config.max_pages and page_num > self.config.max_pages:
                break
                
            text = page.get_text()
            if text.strip():
                content_parts.append(f"\n## Page {page_num}\n\n{text}")
        
        doc.close()
        
        content = '\n'.join(content_parts)
        return content, metadata

    def _convert_docx(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert DOCX to Markdown."""
        from docx import Document
        
        doc = Document(file_path)
        content_parts = []
        
        # Extract core properties
        core_props = doc.core_properties
        metadata = DocumentMetadata(
            title=core_props.title,
            author=core_props.author,
            subject=core_props.subject,
            created=str(core_props.created) if core_props.created else None,
            modified=str(core_props.modified) if core_props.modified else None,
        )
        
        # Convert paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style.name.startswith('Heading'):
                    level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
                    content_parts.append(f"{'#' * int(level)} {para.text}\n")
                else:
                    content_parts.append(para.text + '\n')
        
        # Convert tables
        for table in doc.tables:
            content_parts.append('\n')
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                content_parts.append(row_text + '\n')
            content_parts.append('\n')
        
        content = '\n'.join(content_parts)
        return content, metadata

    def _convert_doc(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert DOC to Markdown (legacy format)."""
        # Try to use antiword or convert to docx first
        try:
            import subprocess
            result = subprocess.run(
                ['antiword', str(file_path)],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                return result.stdout, DocumentMetadata()
        except (subprocess.SubprocessError, FileNotFoundError):
            pass
        
        # Fallback: try using textract
        try:
            import textract
            text = textract.process(str(file_path)).decode('utf-8')
            return text, DocumentMetadata()
        except Exception:
            raise NotImplementedError(
                "DOC format requires 'antiword' or 'textract'. "
                "Please install: apt-get install antiword"
            )

    def _convert_xlsx(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert XLSX to Markdown."""
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        content_parts = []
        
        metadata = DocumentMetadata(
            title=wb.properties.title,
        )
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            content_parts.append(f"\n## Sheet: {sheet_name}\n")
            
            # Convert to markdown table
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else '' for cell in row]
                rows.append(' | '.join(row_data))
            
            if rows:
                # Add header separator
                header_cols = len(rows[0].split(' | '))
                separator = ' | '.join(['---'] * header_cols)
                rows.insert(1, separator)
                content_parts.extend(rows)
            
            content_parts.append('\n')
        
        wb.close()
        content = '\n'.join(content_parts)
        return content, metadata

    def _convert_xls(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert XLS to Markdown."""
        import pandas as pd
        
        content_parts = []
        
        # Read all sheets
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            content_parts.append(f"\n## Sheet: {sheet_name}\n")
            content_parts.append(df.to_markdown(index=False))
            content_parts.append('\n')
        
        content = '\n'.join(content_parts)
        return content, DocumentMetadata()

    def _convert_pptx(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert PPTX to Markdown."""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        content_parts = []
        
        metadata = DocumentMetadata()
        
        for slide_num, slide in enumerate(prs.slides, 1):
            content_parts.append(f"\n## Slide {slide_num}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content_parts.append(shape.text + '\n')
            
            content_parts.append('\n---\n')
        
        content = '\n'.join(content_parts)
        return content, metadata

    def _convert_ppt(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert PPT to Markdown (legacy format)."""
        raise NotImplementedError(
            "PPT format not directly supported. "
            "Please convert to PPTX first."
        )

    def _convert_html(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert HTML to Markdown."""
        html_content = file_path.read_text(encoding='utf-8', errors='ignore')
        
        # Parse with BeautifulSoup to extract metadata
        soup = BeautifulSoup(html_content, 'html.parser')
        
        title_tag = soup.find('title')
        title = title_tag.get_text() if title_tag else None
        
        metadata = DocumentMetadata(title=title)
        
        # Convert to markdown
        content = md(html_content, heading_style='ATX')
        
        return content, metadata

    def _convert_text(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert plain text to Markdown."""
        # Detect encoding
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            detected = chardet.detect(raw_data)
            encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        
        content = file_path.read_text(encoding=encoding, errors='ignore')
        
        metadata = DocumentMetadata(
            encoding=encoding,
            word_count=len(content.split()),
        )
        
        return content, metadata

    def _convert_markdown(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Read Markdown file."""
        content = file_path.read_text(encoding='utf-8', errors='ignore')
        
        metadata = DocumentMetadata(
            word_count=len(content.split()),
        )
        
        return content, metadata

    def _convert_csv(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert CSV to Markdown table."""
        import pandas as pd
        
        # Try different encodings
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        df = None
        
        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if df is None:
            df = pd.read_csv(file_path, encoding='utf-8', errors='ignore')
        
        content = df.to_markdown(index=False)
        
        metadata = DocumentMetadata(
            word_count=len(content.split()),
        )
        
        return content, metadata

    def _convert_json(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert JSON to Markdown."""
        import json
        
        data = json.loads(file_path.read_text(encoding='utf-8'))
        
        content_parts = ["# JSON Document\n\n"]
        content_parts.append("```json\n")
        content_parts.append(json.dumps(data, indent=2, ensure_ascii=False))
        content_parts.append("\n```\n")
        
        content = ''.join(content_parts)
        
        metadata = DocumentMetadata()
        return content, metadata

    def _convert_xml(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert XML to Markdown."""
        content = file_path.read_text(encoding='utf-8', errors='ignore')
        
        # Wrap in code block
        content = f"# XML Document\n\n```xml\n{content}\n```\n"
        
        metadata = DocumentMetadata()
        return content, metadata

    def _convert_rtf(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert RTF to Markdown."""
        from striprtf.striprtf import rtf_to_text
        
        rtf_content = file_path.read_text(encoding='utf-8', errors='ignore')
        content = rtf_to_text(rtf_content)
        
        metadata = DocumentMetadata()
        return content, metadata

    def _convert_odt(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert ODT to Markdown."""
        from odf import opendocument
        from odf.text import P
        
        doc = opendocument.load(file_path)
        paragraphs = doc.spreadsheet.getElementsByType(P)
        
        content = '\n\n'.join([str(p) for p in paragraphs])
        
        metadata = DocumentMetadata()
        return content, metadata

    def _convert_ods(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert ODS to Markdown."""
        import pandas as pd
        
        # Read ODS file
        df = pd.read_excel(file_path, engine='odf')
        content = df.to_markdown(index=False)
        
        metadata = DocumentMetadata()
        return content, metadata

    def _convert_odp(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Convert ODP to Markdown."""
        # For now, treat as binary and extract what we can
        raise NotImplementedError("ODP conversion not yet implemented")
